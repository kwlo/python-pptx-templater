from pptx import Presentation

from pptx_templater.slides import clear


def create(srcpath):
    with open(srcpath, 'rb') as f:
        source = Presentation(f)
        dest = Presentation(f)

        # Clear all slides in destination
        clear(dest)

        return source, dest
